import pptx
import pptx.util
import pptx.dml.color
import pandas as pd
import re
import os
from bs4 import BeautifulSoup


class_name = "sorcerer/wizard"
# sheet_name = "Sorcerer-Wizard"    # not used if full spell list excel is available
# spell_list = ["disguise self"]
spell_list = {"dancing lights", "resistance", "ray of frost", "mage hand", "read magic", "prestidigitation",
              "detect magic", "detect poison", "mage armor", "burning hands", "disguise self", "identify",
              "chastise", "fire breath", "flaming sphere", "color spray", "see invisibility",
              "fireball", "arcane mark", "alarm", "burning arc", "message", "mending", "comprehend languages",
              "silent image", "eagle's splendor", "diminish resistance"}

arcane_class_rgx = "wizard|sorcerer|bard"
divine_class_rgx = "cleric|druid|paladin|ranger|inquisitor|oracle|hunter|shaman|warpriest"


school_colors = {
    "abjuration": [0, 109, 163],
    "conjuration": [206, 163, 0],
    "divination": [164, 198, 204],
    "enchantment": [205, 96, 223],
    "evocation": [148, 39, 22],
    "illusion": [118, 22, 196],
    "necromancy": [16, 53, 14],
    "transmutation": [152, 53, 0],
    "universal": [67, 42, 11]
}


def replace_text(text_frame, new_text, p_idx: int = None):
    if p_idx is None:
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.paragraphs[p_idx]

    if p.runs:
        orig_run = p.runs[0]
        # save formatting
        font = orig_run.font
        font_name = font.name
        font_size = font.size
        font_bold = font.bold
        font_italic = font.italic
        font_underline = font.underline
        try:
            font_color = font.color.rgb
            rgbOk = True
        except AttributeError:
            font_color = font.color.theme_color
            rgbOk = False
        # Clear existing text and apply new text with preserved formatting
        if p_idx is None:
            text_frame.clear()  # Clears all text and formatting
            new_run = text_frame.paragraphs[0].add_run()  # New run in first paragraph
        else:
            p.clear()
            new_run = p.add_run()
        new_run.text = new_text
        # Reapply formatting
        new_run.font.name = font_name
        new_run.font.size = font_size
        new_run.font.bold = font_bold
        new_run.font.italic = font_italic
        new_run.font.underline = font_underline
        if rgbOk:
            new_run.font.color.rgb = font_color
        else:
            new_run.font.color.theme_color = font_color
    else:
        text_frame.text = new_text


def replace_text_by_id(slide, textbox_id, new_text, *args):
    shapes = slide.shapes
    if len(shapes) == 0:
        raise ValueError("Empty shapes!")

    for s in shapes:
        if s.shape_id == textbox_id:
            break

    if s.shape_id != textbox_id:
        raise ValueError(f"Text box with id {textbox_id} was not found")

    replace_text(s.text_frame, new_text, *args)


def replace_text_pretty(text_frame, html_format):
    orig_run = text_frame.paragraphs[0].runs[0]
    # save font name and size
    font = orig_run.font
    font_name = font.name
    font_size = font.size

    # clear other formatting
    font_bold = False
    font_italic = False
    font_underline = False
    text_frame.clear()

    try:
        font_color = font.color.rgb
        rgbOk = True
    except AttributeError:
        font_color = font.color.theme_color
        rgbOk = False

    soup = BeautifulSoup(html_format, 'html.parser')
    current_node = next(next(soup.children).children) # we know that soup contains at least 1 paragraph, which is already added
    visit_stack = []
    for child in reversed(soup.contents[1:]):
        visit_stack.append((child, True))
    for child in reversed(soup.contents[0].contents):
        visit_stack.append((child, True))
    while visit_stack:
        current_node, is_before = visit_stack.pop()
        if current_node.name is None:    # current node is a leaf text
            new_run = text_frame.paragraphs[-1].add_run()
            new_run.text = str(current_node)
            new_run.font.name = font_name
            new_run.font.size = font_size
            new_run.font.bold = font_bold
            new_run.font.italic = font_italic
            new_run.font.underline = font_underline
            if rgbOk:
                new_run.font.color.rgb = font_color
            else:
                new_run.font.color.theme_color = font_color
        elif current_node.name == 'p':
            text_frame.add_paragraph()
            for child in reversed(current_node.contents):
                visit_stack.append((child, True))
        else:
            if current_node.name == 'i':
                font_italic = is_before
            elif current_node.name == 'b':
                font_bold = is_before
            elif current_node.name == 'u':
                font_underline = is_before
            else:
                continue

            if is_before:
                visit_stack.append((current_node, False))
                for child in reversed(current_node.contents):
                    visit_stack.append((child, True))


def populate_metatemplate(class_name, spell_row, slide):
    # insert images
    slide.placeholders[10].insert_picture(f"./assets/backgrounds/{spell_row['school']}.png")
    if re.search(arcane_class_rgx, spell_row['spell_level']):
        slide.placeholders[14].insert_picture("./assets/icons/arcane.png")
    if re.search(divine_class_rgx, spell_row['spell_level']):
        slide.placeholders[15].insert_picture("./assets/icons/divine.png")
    if (not pd.isna(spell_row['spell_resistance'])) and spell_row['spell_resistance'] != 'no':
        slide.placeholders[16].insert_picture("./assets/icons/shield.png")

    # create name and components
    name_tf = slide.placeholders[11].text_frame
    name_tf.clear()
    sp_name_run = name_tf.paragraphs[0].add_run()
    sp_name_run.text = spell_row['name']
    sp_name_run.font.name = 'Amasis MT Pro Black'
    sp_name_run.font.size = pptx.util.Pt(28)
    sp_name_run.font.color.rgb = pptx.dml.color.RGBColor(*(school_colors[spell_row['school']]))

    try:
        comp_str = re.search(".*\\(", spell_row['components']).group()[:-1].strip()
    except AttributeError:
        comp_str = spell_row['components']
    name_tf.add_paragraph()
    ing_run = name_tf.paragraphs[1].add_run()
    ing_run.text = comp_str
    ing_run.font.name = 'Amasis MT Pro Medium'
    ing_run.font.size = pptx.util.Pt(18)
    ing_run.font.color.rgb = pptx.dml.color.RGBColor(*(school_colors[spell_row['school']]))
    name_tf.paragraphs[1].space_before = pptx.util.Pt(3)

    try:
        det_str = re.search("(\\(.*\\))", spell_row['components']).group()[1:-1].strip()
        name_tf.add_paragraph()
        ing_run = name_tf.paragraphs[2].add_run()
        ing_run.text = det_str
        ing_run.font.name = 'Amasis MT Pro'
        ing_run.font.size = pptx.util.Pt(14)
        ing_run.font.color.rgb = pptx.dml.color.RGBColor(*(school_colors[spell_row['school']]))
        name_tf.paragraphs[2].space_before = pptx.util.Pt(1)
    except AttributeError:
        pass

    # add level
    lvl = re.search(f"{class_name} [0-9]+", spell_row['spell_level']).group()[-1].strip()
    replace_text(slide.placeholders[13].text_frame, lvl)

    # build description table
    table = slide.placeholders[12].table

    # determine data
    datas = [["Time", spell_row['casting_time']], ["Duration", spell_row['duration']], ["Range", spell_row['range']]]
    if (not pd.isna(spell_row['targets'])) and (not pd.isna(spell_row['area'])):
        if spell_row['targets'] == spell_row['area']:
            datas.append(["Area/Target", spell_row['targets']])
        else:
            raise Exception("What to do if both targets and area is defined?")
    if not pd.isna(spell_row['targets']):
        datas.append(["Targets", spell_row['targets']])
    elif not pd.isna(spell_row['area']):
        datas.append(["Area", spell_row['area']])

    if not (spell_row['saving_throw'] == 'none' or pd.isna(spell_row['saving_throw'])):
        datas.append(["Save", spell_row['saving_throw']])

    if not pd.isna(spell_row['effect']):
        datas.append(["Effect", spell_row['effect']])

    if len(datas) % 2 == 1:
        datas.append(["", ""])

    if len(datas) < 5:
        table._tbl.remove(table.rows[0]._tr)
        desc_row_idx = 2
    else:
        desc_row_idx = 3

    for idx, d in enumerate(datas):
        r_idx = int(idx/2)
        c_idx = (idx % 2) * 2

        replace_text(table.cell(r_idx,c_idx).text_frame, d[0])
        replace_text(table.cell(r_idx, c_idx+1).text_frame, d[1])

    replace_text_pretty(table.cell(desc_row_idx, 0).text_frame, spell_row['description_formatted'])


if __name__ == "__main__":
    print("Loading spell database")
    # df = pd.read_excel("spells_by_class.xlsx", sheet_name)
    df = pd.read_excel("spell_full.xlsx")
    print("Loading done")

    spell_list = pd.Series(list(spell_list))
    spell_list = spell_list.str.lower()

    lower_names = df['name'].str.lower()
    selected_spells = df.loc[lower_names.isin(spell_list)]
    found_num = selected_spells.shape[0]
    print(f"{found_num}/{spell_list.size} spells have been found.")
    print("Initiate export")
    if not os.path.isdir('./output'):
        print("Output folder not found, create it")
        os.mkdir('./output')

    if not os.path.isdir('./output/cards'):
        print("Output folder not found, create it")
        os.mkdir('./output/cards')

    for cnt, idx_row in enumerate(selected_spells.iterrows()):
        spell_row = idx_row[1]
        presentation = pptx.Presentation(f"assets/template.pptx")
        populate_metatemplate(class_name, spell_row, presentation.slides[0])
        presentation.save(f"output/cards/{spell_row['name']}.pptx")
        print(f"{cnt+1}/{found_num} done")

    print("Exporting done")
    not_found_spells = spell_list[~spell_list.isin(lower_names)]
    if not not_found_spells.empty:
        print("The following spells have not been found:")
        print(not_found_spells)
